from flask import Flask, render_template, request, jsonify, send_file, send_from_directory
import json
import io
import os
import base64
import uuid
from datetime import datetime
import openai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.text import MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
import re
import math
import asyncio
import concurrent.futures
from concurrent.futures import ThreadPoolExecutor

app = Flask(__name__)

# Initialize OpenAI client (users should set OPENAI_API_KEY environment variable)
openai_client = openai.OpenAI(api_key=os.getenv('OPENAI_API_KEY'))

# Create images directory if it doesn't exist
IMAGES_DIR = os.path.join(os.path.dirname(__file__), 'static', 'generated_images')
os.makedirs(IMAGES_DIR, exist_ok=True)

@app.route('/static/generated_images/<filename>')
def serve_generated_image(filename):
    """Serve generated images"""
    return send_from_directory(IMAGES_DIR, filename)

def parse_bullet_points(content):
    """Parse content and extract bullet points - each line becomes a bullet point"""
    if not content:
        return []
    
    lines = content.strip().split('\n')
    bullet_points = []
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        # Remove any bullet symbols that might be present (including hyphens)
        if line.startswith(('- ', '• ', '* ', '· ', '○ ', '▪ ', '‣ ')):
            bullet_text = line[2:].strip()
            if bullet_text:
                bullet_points.append(bullet_text)
        elif line.startswith('-'):
            # Handle cases where there's just a dash without space
            bullet_text = line[1:].strip()
            if bullet_text:
                bullet_points.append(bullet_text)
        elif line and not any(line.lower().startswith(word) for word in ['slide', 'title:', 'content:']):
            # Include non-empty lines that aren't headers - each line becomes a bullet
            bullet_points.append(line)
    
    return bullet_points


@app.route('/')
def index():
    return render_template('index.html')

@app.route('/generate_code', methods=['POST'])
def generate_code():
    layout_data = request.json
    
    # Generate python-pptx code
    code = generate_pptx_code(layout_data)
    
    # Create downloadable file
    file_content = io.StringIO(code)
    file_content.seek(0)
    
    return jsonify({
        'code': code,
        'filename': f'powerpoint_layout_{datetime.now().strftime("%Y%m%d_%H%M%S")}.py'
    })

@app.route('/download_code', methods=['POST'])
def download_code():
    code = request.json.get('code')
    filename = request.json.get('filename', 'powerpoint_layout.py')
    
    # Create file in memory
    file_buffer = io.BytesIO()
    file_buffer.write(code.encode('utf-8'))
    file_buffer.seek(0)
    
    return send_file(
        file_buffer,
        as_attachment=True,
        download_name=filename,
        mimetype='text/plain'
    )

def generate_pptx_code(layout_data):
    """Generate python-pptx code from layout data"""
    
    code_lines = [
        "from pptx import Presentation",
        "from pptx.util import Inches, Pt",
        "from pptx.enum.text import PP_ALIGN",
        "from pptx.enum.text import MSO_ANCHOR",
        "from pptx.dml.color import RGBColor",
        "",
        "# Create presentation",
        "pres = Presentation()",
        ""
    ]
    
    slide_type = layout_data.get('slide_type', 'content')
    elements = layout_data.get('elements', [])
    
    if slide_type == 'title':
        code_lines.append("# Add title slide")
        code_lines.append("slide_layout = pres.slide_layouts[6]  # Blank layout")
    else:
        code_lines.append("# Add content slide") 
        code_lines.append("slide_layout = pres.slide_layouts[6]  # Blank layout")
    
    code_lines.append("slide = pres.slides.add_slide(slide_layout)")
    code_lines.append("")
    
    # Add elements
    for i, element in enumerate(elements):
        if element['type'] in ['textbox', 'title']:
            code_lines.extend(generate_textbox_code(element, i))
        elif element['type'] == 'image':
            code_lines.extend(generate_image_code(element, i))
    
    code_lines.extend([
        "",
        "# Save presentation",
        "pres.save('generated_presentation.pptx')",
        "print('Presentation saved as generated_presentation.pptx')"
    ])
    
    return "\n".join(code_lines)

def generate_textbox_code(element, index):
    """Generate code for text box element"""
    
    left = element['left']
    top = element['top'] 
    width = element['width']
    height = element['height']
    font_size = element.get('font_size', 18)
    font_name = element.get('font_name', 'Arial')
    list_type = element.get('list_type', 'none')
    content = element.get('content', 'Sample text')
    
    lines = [
        f"# Add {'title' if element['type'] == 'title' else 'text'} box {index + 1}",
        f"textbox_{index} = slide.shapes.add_textbox(",
        f"    Inches({left}), Inches({top}), Inches({width}), Inches({height})",
        f")",
        f"text_frame_{index} = textbox_{index}.text_frame",
        f"text_frame_{index}.clear()  # Clear default paragraph",
        f"text_frame_{index}.word_wrap = True",
        ""
    ]
    
    # Check if content has bullet points
    if list_type == 'bullet' and ('\\n' in content or '\n' in content):
        bullet_points = parse_bullet_points(content)
        lines.extend([
            f"# Add bullet points",
            f"bullet_items = {repr(bullet_points)}",
            f"for i, item in enumerate(bullet_items):",
            f"    if item.strip():",
            f"        if i == 0:",
            f"            p = text_frame_{index}.paragraphs[0]",
            f"        else:",
            f"            p = text_frame_{index}.add_paragraph()",
            f"        p.text = f'• {{item.strip()}}'",
            f"        p.level = 0  # First level bullet",
            f"        # Format the paragraph",
            f"        for run in p.runs:",
            f"            font = run.font",
            f"            font.name = '{font_name}'",
            f"            font.size = Pt({font_size})",
        ])
    else:
        # Simple text without bullets
        lines.extend([
            f"p = text_frame_{index}.paragraphs[0]",
            f"p.text = '{content}'",
            f"",
            f"# Format text",
            f"for run in p.runs:",
            f"    font = run.font",
            f"    font.name = '{font_name}'",
            f"    font.size = Pt({font_size})",
        ])
        
        # Add alignment for title
        if element['type'] == 'title':
            lines.insert(-5, "p.alignment = PP_ALIGN.CENTER")
    
    lines.append("")
    
    return lines

def generate_image_code(element, index):
    """Generate code for image element"""
    
    left = element['left']
    top = element['top']
    width = element['width'] 
    height = element['height']
    
    lines = [
        f"# Add image {index + 1}",
        f"# Note: Replace 'path_to_image.jpg' with your actual image path",
        f"try:",
        f"    # Try to add actual image",
        f"    image_{index} = slide.shapes.add_picture(",
        f"        'path_to_image.jpg',",
        f"        Inches({left}), Inches({top}), Inches({width}), Inches({height})",
        f"    )",
        f"except:",
        f"    # If image file not found, add placeholder rectangle",
        f"    from pptx.enum.shapes import MSO_SHAPE",
        f"    rectangle_{index} = slide.shapes.add_shape(",
        f"        MSO_SHAPE.RECTANGLE,",
        f"        Inches({left}), Inches({top}), Inches({width}), Inches({height})",
        f"    )",
        f"    # Style as image placeholder",
        f"    fill = rectangle_{index}.fill",
        f"    fill.solid()",
        f"    fill.fore_color.rgb = RGBColor(240, 240, 240)",
        f"    line = rectangle_{index}.line",
        f"    line.color.rgb = RGBColor(169, 169, 169)",
        f"    line.width = Pt(1)",
        f"    # Add placeholder text",
        f"    if rectangle_{index}.has_text_frame:",
        f"        text_frame = rectangle_{index}.text_frame",
        f"        text_frame.clear()",
        f"        p = text_frame.paragraphs[0]",
        f"        p.text = '[INSERT IMAGE HERE]'",
        f"        p.alignment = PP_ALIGN.CENTER",
        f"        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE",
        ""
    ]
    
    return lines

def generate_simple_image_prompt(slide_title, slide_content=None):
    """Generate a simple, direct image prompt for a slide"""
    try:
        # Extract first bullet point if content exists
        first_bullet = ""
        if slide_content:
            bullet_points = parse_bullet_points(slide_content)
            if bullet_points:
                first_bullet = bullet_points[0]
        
        # Simple, direct prompt generation
        response = openai_client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "Create a simple, direct image prompt for an educational slide. The image should clearly illustrate the main topic. Be literal and specific. No artistic interpretations."},
                {"role": "user", "content": f"Create an image prompt for a slide titled '{slide_title}'.\n{f'The slide discusses: {first_bullet}' if first_bullet else ''}\n\nGenerate a simple, direct prompt that shows exactly what the slide is about. Example: If the slide is about 'Baboon Social Structure', the prompt should be 'A group of baboons grooming each other'."}
            ],
            max_tokens=50,
            temperature=0.3  # Lower temperature for more consistent results
        )
        
        return response.choices[0].message.content.strip()
    except:
        # Simple fallback - no additional API calls
        return f"An educational illustration showing {slide_title.lower()}"

@app.route('/generate_draft', methods=['POST'])
def generate_draft():
    """Generate a draft outline of slide titles based on user topic"""
    data = request.json
    topic = data.get('topic', '')
    
    if not topic:
        return jsonify({'error': 'Topic is required'}), 400
    
    try:
        # Generate slide outline using OpenAI
        response = openai_client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are a presentation expert. Create a clear, logical outline for a PowerPoint presentation. The first line should be the main presentation title (what goes on the title slide). The following lines should be content slide titles. Return only slide titles, one per line, without numbering, bullet points, or prefixes like 'Title:' or 'Slide 1:'. Provide 5-9 total lines."},
                {"role": "user", "content": f"Create a presentation outline about: {topic}\n\nExample format:\nThe Complete Guide to Solar Energy\nIntroduction to Solar Technology\nTypes of Solar Panels\nInstallation Process\nCost and Benefits\nMaintenance and Care"}
            ],
            max_tokens=300,
            temperature=0.7
        )
        
        slides_text = response.choices[0].message.content.strip()
        slides = [line.strip() for line in slides_text.split('\n') if line.strip()]
        
        # Format slides with types and clean up titles
        formatted_slides = []
        
        # First slide is always the title slide using the first line from OpenAI
        if slides:
            title_slide_title = slides[0]
            # Remove common prefixes
            title_slide_title = re.sub(r'^Slide\s*\d+\s*:\s*', '', title_slide_title, flags=re.IGNORECASE)
            title_slide_title = re.sub(r'^\d+\.\s*', '', title_slide_title)
            
            formatted_slides.append({
                'id': 0,
                'title': title_slide_title,
                'type': 'title',
                'content': ''  # Title slides typically have no content or just subtitle
            })
            
            # Remaining slides are content slides
            for i, title in enumerate(slides[1:], 1):
                # Remove common prefixes
                cleaned_title = re.sub(r'^Slide\s*\d+\s*:\s*', '', title, flags=re.IGNORECASE)
                cleaned_title = re.sub(r'^\d+\.\s*', '', cleaned_title)
                
                formatted_slides.append({
                    'id': i,
                    'title': cleaned_title,
                    'type': 'content',
                    'content': 'Content will be generated after approval'
                })
        
        return jsonify({
            'slides': formatted_slides,
            'topic': topic
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/generate_content', methods=['POST'])
def generate_content():
    """Generate full content for approved slide outline"""
    data = request.json
    slides = data.get('slides', [])
    topic = data.get('topic', '')
    content_layout = data.get('content_layout', {})
    
    if not slides:
        return jsonify({'error': 'Slides are required'}), 400
    
    
    try:
        # Generate content for each slide
        for slide in slides:
            if slide['type'] == 'content' and not slide.get('content_generated', False):
                response = openai_client.chat.completions.create(
                    model="gpt-3.5-turbo",
                    messages=[
                        {"role": "system", "content": "You are a presentation content writer. Create bullet points for PowerPoint slides. Provide each point as a separate line with NO bullet symbols, NO dashes, NO prefixes - just plain text. Each line will automatically become a bullet point in the presentation. Be concise and impactful."},
                        {"role": "user", "content": f"Create content for this slide about '{topic}':\nSlide title: {slide['title']}\n\nProvide 3-5 bullet points. Each point should be on its own line with no bullet symbols."}
                    ],
                    max_tokens=300,
                    temperature=0.7
                )
                
                slide['content'] = response.choices[0].message.content.strip()
                slide['content_generated'] = True
        
        return jsonify({'slides': slides})
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

def generate_single_image(slide_title, slide_content, custom_prompt=None):
    """Generate a single image using gpt-image-1"""
    try:
        # Use custom prompt if provided, otherwise generate simple prompt
        if custom_prompt:
            image_prompt = custom_prompt
        else:
            # Use the new simple prompt generation
            image_prompt = generate_simple_image_prompt(slide_title, slide_content)
        
        # Log the prompt for debugging
        print(f"Generating image for '{slide_title}' with prompt: {image_prompt}")
        
        # Simplified parameters for gpt-image-1
        response = openai_client.images.generate(
            model="gpt-image-1",
            prompt=image_prompt,
            size="1024x1024",
            quality="auto",  # Changed to "auto" - a supported value
            # Removed background parameter
            moderation="low"
        )
        
        # gpt-image-1 returns base64 data, save as file and return URL
        image_base64 = response.data[0].b64_json
        
        # Generate unique filename
        image_filename = f"generated_{uuid.uuid4().hex[:8]}_{int(datetime.now().timestamp())}.png"
        image_path = os.path.join(IMAGES_DIR, image_filename)
        
        # Save base64 image to file
        image_bytes = base64.b64decode(image_base64)
        with open(image_path, 'wb') as f:
            f.write(image_bytes)
        
        # Return URL that can be served by Flask
        image_url = f"/static/generated_images/{image_filename}"
        
        # Instead of using OpenAI for captions, use simple logic:
        caption = f"{slide_title} illustration"
        
        return {
            'image_url': image_url,
            'caption': caption,
            'prompt_used': image_prompt,
            'success': True
        }
        
    except Exception as e:
        print(f"Image generation error: {str(e)}")
        return {
            'error': str(e),
            'success': False
        }

@app.route('/generate_image', methods=['POST'])
def generate_image():
    """Generate an image for a slide using OpenAI gpt-image-1"""
    data = request.json
    slide_title = data.get('title', '')
    slide_content = data.get('content', '')
    custom_prompt = data.get('custom_prompt', None)
    
    if not slide_title:
        return jsonify({'error': 'Slide title is required'}), 400
    
    result = generate_single_image(slide_title, slide_content, custom_prompt)
    
    if result['success']:
        return jsonify({
            'image_url': result['image_url'],
            'caption': result['caption'],
            'prompt_used': result['prompt_used']
        })
    else:
        return jsonify({'error': result['error']}), 500

@app.route('/generate_image_prompt', methods=['POST'])
def generate_image_prompt():
    """Generate a simple image prompt for a slide"""
    data = request.json
    slide_title = data.get('title', '')
    slide_content = data.get('content', '')
    
    if not slide_title:
        return jsonify({'error': 'Slide title is required'}), 400
    
    prompt = generate_simple_image_prompt(slide_title, slide_content)
    
    return jsonify({'prompt': prompt})

@app.route('/generate_images_bulk', methods=['POST'])
def generate_images_bulk():
    """Generate images for multiple slides concurrently using gpt-image-1"""
    data = request.json
    slides = data.get('slides', [])
    
    if not slides:
        return jsonify({'error': 'Slides are required'}), 400
    
    try:
        # Use ThreadPoolExecutor for concurrent image generation
        with ThreadPoolExecutor(max_workers=5) as executor:
            # Submit all image generation tasks
            future_to_index = {}
            for i, slide in enumerate(slides):
                if not slide.get('generated_image'):  # Only generate if no image exists
                    # Use the suggested prompt if available, otherwise generate
                    custom_prompt = slide.get('suggested_image_prompt', None)
                    future = executor.submit(generate_single_image, slide.get('title', ''), slide.get('content', ''), custom_prompt)
                    future_to_index[future] = i
            
            # Collect results as they complete
            results = {}
            for future in concurrent.futures.as_completed(future_to_index):
                slide_index = future_to_index[future]
                result = future.result()
                results[slide_index] = result
        
        # Update slides with generated images
        for slide_index, result in results.items():
            if result['success']:
                slides[slide_index]['generated_image'] = result['image_url']
                slides[slide_index]['image_caption'] = result['caption']
            else:
                slides[slide_index]['image_error'] = result['error']
                print(f"Slide {slide_index} error: {result['error']}")
        
        return jsonify({
            'slides': slides,
            'generated_count': len([r for r in results.values() if r['success']]),
            'error_count': len([r for r in results.values() if not r['success']])
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500

def _add_image_placeholder(slide, element):
    """Helper function to add image placeholder rectangle"""
    rectangle = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(element['left']), 
        Inches(element['top']), 
        Inches(element['width']), 
        Inches(element['height'])
    )
    
    # Style the rectangle to look like an image placeholder
    fill = rectangle.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 240, 240)  # Light gray background
    
    line = rectangle.line
    line.color.rgb = RGBColor(169, 169, 169)  # Gray border
    line.width = Pt(1)
    
    # Add text to indicate this is an image placeholder
    if rectangle.has_text_frame:
        text_frame = rectangle.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.text = "[INSERT IMAGE HERE]"
        p.alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Format the placeholder text
        if p.runs:
            for run in p.runs:
                font = run.font
                font.name = 'Calibri'
                font.size = Pt(14)
                font.color.rgb = RGBColor(128, 128, 128)  # Gray text

@app.route('/create_presentation', methods=['POST'])
def create_presentation():
    """Create final PPTX file with generated content and custom layouts"""
    data = request.json
    slides_data = data.get('slides', [])
    title_layout = data.get('title_layout', {})
    content_layout = data.get('content_layout', {})
    
    if not slides_data:
        return jsonify({'error': 'Slides data is required'}), 400
    
    # Debug: print the received layouts
    print("Title layout received:", title_layout)
    print("Content layout received:", content_layout)
    
    try:
        # Create presentation with 16:9 aspect ratio
        pres = Presentation()
        
        # Set slide size to 16:9 widescreen format
        pres.slide_width = Inches(13.333)  # 16:9 ratio width
        pres.slide_height = Inches(7.5)    # 16:9 ratio height
        
        for slide_data in slides_data:
            slide_type = slide_data['type']
            
            # Use the appropriate layout based on slide type
            if slide_type == 'title':
                layout_config = title_layout
            else:
                layout_config = content_layout
            
            print(f"Using layout for {slide_type} slide:", layout_config)
            
            # Add blank slide
            slide_layout = pres.slide_layouts[6]  # Blank layout
            slide = pres.slides.add_slide(slide_layout)
            
            # Add elements based on layout configuration
            elements = layout_config.get('elements', [])
            print(f"Processing {len(elements)} elements for slide: {slide_data['title']}")
            
            for element in elements:
                print(f"Processing element: {element}")
                
                if element['type'] in ['textbox', 'title']:
                    # Add text box with EXACT positioning from the designer
                    textbox = slide.shapes.add_textbox(
                        Inches(element['left']), 
                        Inches(element['top']), 
                        Inches(element['width']), 
                        Inches(element['height'])
                    )
                    text_frame = textbox.text_frame
                    text_frame.clear()  # Clear default content
                    
                    # Configure text frame properties
                    text_frame.word_wrap = True
                    # Don't use auto-sizing to preserve exact font sizes
                    text_frame.auto_size = MSO_AUTO_SIZE.NONE
                    
                    # Determine what content to use
                    if element['type'] == 'title':
                        # For title elements, always use the slide title
                        content_text = slide_data['title']
                        p = text_frame.paragraphs[0]
                        p.text = content_text
                        p.alignment = PP_ALIGN.CENTER
                        
                        # Apply title formatting from the designer - force run creation
                        if not p.runs:
                            text_content = p.text
                            p.text = text_content
                            
                        for run in p.runs:
                            font = run.font
                            font.name = element.get('font_name', 'Calibri')
                            font.size = Pt(element.get('font_size', 28))
                            font.bold = True
                    else:
                        # For textbox elements, use slide content
                        content_text = slide_data.get('content', '')
                        
                        if slide_type == 'title':
                            # For title slides, textbox shows subtitle/description
                            p = text_frame.paragraphs[0]
                            p.text = content_text if content_text else ""
                            p.alignment = PP_ALIGN.CENTER
                            
                            # Force run creation and apply formatting
                            if not p.runs:
                                text_content = p.text
                                p.text = text_content
                                
                            for run in p.runs:
                                font = run.font
                                font.name = element.get('font_name', 'Calibri')
                                font.size = Pt(element.get('font_size', 18))
                        else:
                            # For content slides, parse and add bullet points if list_type is bullet
                            if element.get('list_type', 'none') == 'bullet':
                                bullet_points = parse_bullet_points(content_text)
                                print(f"Parsed bullet points: {bullet_points}")
                                
                                if bullet_points:
                                    # Add bullet points - each line becomes a bullet
                                    for i, bullet_text in enumerate(bullet_points):
                                        if i == 0:
                                            p = text_frame.paragraphs[0]
                                        else:
                                            p = text_frame.add_paragraph()
                                        
                                        # Add bullet character manually for blank layouts
                                        p.text = f"• {bullet_text.strip()}"
                                        p.level = 0  # First level bullet
                                        
                                        # Apply formatting from the designer - force run creation
                                        if not p.runs:
                                            # Force run creation by setting text again
                                            text_content = p.text
                                            p.text = text_content
                                        
                                        # Apply to all runs
                                        for run in p.runs:
                                            run_font = run.font
                                            run_font.name = element.get('font_name', 'Calibri')
                                            run_font.size = Pt(element.get('font_size', 18))
                                else:
                                    # No bullet points, just add the text
                                    p = text_frame.paragraphs[0]
                                    p.text = content_text
                                    
                                    # Force run creation and apply formatting
                                    if not p.runs:
                                        text_content = p.text
                                        p.text = text_content
                                    
                                    for run in p.runs:
                                        font = run.font
                                        font.name = element.get('font_name', 'Calibri')
                                        font.size = Pt(element.get('font_size', 18))
                            else:
                                # Plain text without bullets
                                p = text_frame.paragraphs[0]
                                p.text = content_text
                                
                                # Force run creation and apply formatting
                                if not p.runs:
                                    text_content = p.text
                                    p.text = text_content
                                
                                for run in p.runs:
                                    font = run.font
                                    font.name = element.get('font_name', 'Calibri')
                                    font.size = Pt(element.get('font_size', 18))
                                        
                elif element['type'] == 'image':
                    # Check if the slide has a generated image
                    generated_image_url = slide_data.get('generated_image', '')
                    
                    if generated_image_url and generated_image_url.startswith('/static/generated_images/'):
                        # Extract filename from URL and construct full path
                        image_filename = generated_image_url.replace('/static/generated_images/', '')
                        image_path = os.path.join(IMAGES_DIR, image_filename)
                        
                        # Check if the image file exists
                        if os.path.exists(image_path):
                            try:
                                # Add the actual generated image
                                picture = slide.shapes.add_picture(
                                    image_path,
                                    Inches(element['left']), 
                                    Inches(element['top']), 
                                    Inches(element['width']), 
                                    Inches(element['height'])
                                )
                                print(f"Successfully added image: {image_filename}")
                            except Exception as e:
                                print(f"Error adding image {image_filename}: {str(e)}")
                                # Fall back to placeholder if image loading fails
                                _add_image_placeholder(slide, element)
                        else:
                            print(f"Image file not found: {image_path}")
                            # Fall back to placeholder if file doesn't exist
                            _add_image_placeholder(slide, element)
                    else:
                        print(f"No generated image found for slide: {slide_data['title']}")
                        # Add placeholder if no generated image
                        _add_image_placeholder(slide, element)
        
        # Save presentation to memory
        file_buffer = io.BytesIO()
        pres.save(file_buffer)
        file_buffer.seek(0)
        
        filename = f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        
        return send_file(
            file_buffer,
            as_attachment=True,
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
        
    except Exception as e:
        print("Error creating presentation:", str(e))
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    app.run(debug=True)
